"""
Phase 7: PowerPoint Generation Nodes (Steps 21-22)

This module contains nodes for generating PowerPoint presentations:
- Step 21: generate_powerpoint_node - Create PowerPoint with charts from tables
- Step 22: (Future) HTML Dashboard generation

Chart Type Selection:
The module intelligently selects chart types based on table dimensions and semantics:
- 2x2 tables: Clustered column (simple comparison)
- >5 rows: Horizontal bar (better for many categories)
- >4 columns: 100% stacked column (part-to-whole relationships)
- Time series data: Line chart
- Ranking data: Horizontal bar
- Default: Clustered column
"""

import logging
from typing import Dict, Any, Optional, Tuple, List
from enum import Enum
from datetime import datetime
from pathlib import Path

from agent.state import WorkflowState, STEP_21_GENERATE_POWERPOINT
from agent.config import DEFAULT_CONFIG
from agent.styling import (
    get_market_research_color,
    get_chart_color,
    hex_to_rgb,
    TITLE_FONT,
    HEADING_FONT,
    BODY_FONT,
    FOOTER_FONT,
    TABLE_HEADER_FONT,
    TABLE_CELL_FONT,
)
from agent.utils.tracing import trace_node

logger = logging.getLogger(__name__)


# =============================================================================
# Chart Type Constants
# =============================================================================

class ChartType(str, Enum):
    """
    Chart type enumeration for PowerPoint generation.

    Each chart type maps to a python-pptx XL_CHART_TYPE enum value and is
    optimized for specific table structures and data patterns.

    Attributes:
        CLUSTERED_COLUMN: Standard vertical bar chart for simple comparisons
        HORIZONTAL_BAR: Horizontal bar chart for many row categories
        STACKED_COLUMN_100: 100% stacked column for part-to-whole with many columns
        LINE: Line chart for time series or ordered categories
    """
    CLUSTERED_COLUMN = "clustered_column"
    HORIZONTAL_BAR = "horizontal_bar"
    STACKED_COLUMN_100 = "stacked_column_100"
    LINE = "line"


# =============================================================================
# Chart Type Documentation
# =============================================================================

CHART_TYPE_DOCUMENTATION = {
    ChartType.CLUSTERED_COLUMN: {
        "xl_chart_type": "COLUMN_CLUSTERED",
        "description": "Standard vertical bar chart with grouped bars",
        "best_for": [
            "Simple comparisons between categories",
            "Small tables (2-5 rows, 2-4 columns)",
            "Side-by-side value comparisons",
            "2x2 tables (two categories, two variables)"
        ],
        "readability": "Excellent for small datasets, becomes cluttered with many categories",
        "examples": [
            "Gender vs Satisfaction (2x2)",
            "Age Group vs Product Preference (3-5 rows, 2-3 columns)"
        ]
    },
    ChartType.HORIZONTAL_BAR: {
        "xl_chart_type": "BAR_CLUSTERED",
        "description": "Horizontal bar chart with categories on Y-axis",
        "best_for": [
            "Tables with many row categories (>5 rows)",
            "Long category labels that need horizontal space",
            "Ranking or ordering data",
            "Single row with many columns"
        ],
        "readability": "Best for tables with many rows; easier to read long labels",
        "examples": [
            "Region vs Satisfaction (10+ regions)",
            "Brand preference ranking (many brands)"
        ]
    },
    ChartType.STACKED_COLUMN_100: {
        "xl_chart_type": "COLUMN_STACKED_100",
        "description": "100% stacked column chart showing part-to-whole relationships",
        "best_for": [
            "Tables with many columns (>4 columns)",
            "Part-to-whole analysis",
            "Percentage distributions across categories",
            "Comparing relative proportions"
        ],
        "readability": "Good for showing relative proportions, not absolute values",
        "examples": [
            "Department vs Response Category (5+ response options)",
            "Region vs Agreement Level (Likert scale with many options)"
        ]
    },
    ChartType.LINE: {
        "xl_chart_type": "LINE",
        "description": "Line chart for ordered or sequential data",
        "best_for": [
            "Time series data (ordered by time)",
            "Ordered categories (e.g., low, medium, high)",
            "Showing trends across ordered dimensions"
        ],
        "readability": "Excellent for trends, misleading for unordered categories",
        "examples": [
            "Year vs Satisfaction (time series)",
            "Income Bracket vs Response (ordered categories)"
        ]
    }
}


# =============================================================================
# Semantic Hints
# =============================================================================

class SemanticHint(str, Enum):
    """
    Semantic hints to guide chart type selection beyond dimensions.

    These hints can be provided by LLM analysis or manual annotation to
    override the dimension-based selection logic.

    Attributes:
        TIME_SERIES: Data is ordered by time; use line chart
        RANKING: Data represents rankings; use horizontal bar
        PART_TO_WHOLE: Data shows proportions; use stacked column
        COMPARISON: Data is for comparison; use clustered column
    """
    TIME_SERIES = "time_series"
    RANKING = "ranking"
    PART_TO_WHOLE = "part_to_whole"
    COMPARISON = "comparison"


# =============================================================================
# Chart Type Selection Function
# =============================================================================

def select_chart_type(
    table_data: Dict[str, Any],
    semantic_hint: Optional[str] = None
) -> str:
    """
    Select appropriate chart type based on table dimensions and semantic hints.

    This function implements a decision tree that considers:
    1. Semantic hints (if provided, these take precedence)
    2. Table dimensions (rows x columns)
    3. Readability best practices

    Selection Logic:
    ----------------
    Priority 1: Semantic hints override dimension-based selection
    - time_series → LINE chart (for ordered/time-based data)
    - ranking → HORIZONTAL_BAR (for ranked categories)

    Priority 2: Special dimension cases
    - 2x2 table → CLUSTERED_COLUMN (simple comparison)
    - Single row → HORIZONTAL_BAR (better label readability)
    - Single column → CLUSTERED_COLUMN (standard comparison)

    Priority 3: Dimension-based selection
    - >5 rows → HORIZONTAL_BAR (many categories)
    - >4 columns → STACKED_COLUMN_100 (many columns = part-to-whole)

    Priority 4: Default
    - CLUSTERED_COLUMN (general-purpose chart)

    Args:
        table_data: Table data dictionary containing:
            - data: Dict with row_labels and column_labels
            - table_name: Optional table identifier
        semantic_hint: Optional semantic hint from SemanticHint enum:
            - "time_series": Use line chart
            - "ranking": Use horizontal bar chart
            - "part_to_whole": Use 100% stacked column
            - "comparison": Use clustered column
            - None: Auto-select based on dimensions

    Returns:
        Chart type string from ChartType enum values

    Raises:
        ValueError: If table_data is malformed or missing required fields

    Examples:
        >>> table_2x2 = {
        ...     "data": {
        ...         "row_labels": ["Male", "Female"],
        ...         "column_labels": ["Yes", "No"]
        ...     }
        ... }
        >>> select_chart_type(table_2x2)
        'clustered_column'

        >>> table_many_rows = {
        ...     "data": {
        ...         "row_labels": [f"Region{i}" for i in range(10)],
        ...         "column_labels": ["Yes", "No"]
        ...     }
        ... }
        >>> select_chart_type(table_many_rows)
        'horizontal_bar'

        >>> table_time_series = {
        ...     "data": {
        ...         "row_labels": ["2020", "2021", "2022"],
        ...         "column_labels": ["Yes", "No"]
        ...     }
        ... }
        >>> select_chart_type(table_time_series, semantic_hint="time_series")
        'line'
    """
    # ==========================================================================
    # Validate input
    # ==========================================================================
    if not table_data:
        raise ValueError("table_data cannot be None or empty")

    data_section = table_data.get("data")
    if not data_section:
        raise ValueError("table_data must contain 'data' section")

    row_labels = data_section.get("row_labels")
    column_labels = data_section.get("column_labels")

    if row_labels is None or column_labels is None:
        raise ValueError(
            "table_data.data must contain 'row_labels' and 'column_labels'"
        )

    # ==========================================================================
    # Extract dimensions
    # ==========================================================================
    n_rows = len(row_labels)
    n_cols = len(column_labels)

    table_name = table_data.get("table_name", "unknown")
    logger.debug(
        f"select_chart_type: {table_name} - dimensions: {n_rows} rows x {n_cols} cols, "
        f"hint: {semantic_hint}"
    )

    # ==========================================================================
    # Priority 1: Semantic hints (override dimension-based selection)
    # ==========================================================================
    if semantic_hint:
        if semantic_hint == SemanticHint.TIME_SERIES:
            logger.debug(f"  → LINE chart (semantic hint: time_series)")
            return ChartType.LINE.value

        if semantic_hint == SemanticHint.RANKING:
            logger.debug(f"  → HORIZONTAL_BAR chart (semantic hint: ranking)")
            return ChartType.HORIZONTAL_BAR.value

        if semantic_hint == SemanticHint.PART_TO_WHOLE:
            logger.debug(f"  → STACKED_COLUMN_100 chart (semantic hint: part_to_whole)")
            return ChartType.STACKED_COLUMN_100.value

        if semantic_hint == SemanticHint.COMPARISON:
            logger.debug(f"  → CLUSTERED_COLUMN chart (semantic hint: comparison)")
            return ChartType.CLUSTERED_COLUMN.value

    # ==========================================================================
    # Priority 2: Special dimension cases
    # ==========================================================================

    # 2x2 table → Clustered column (simple comparison)
    if n_rows == 2 and n_cols == 2:
        logger.debug(f"  → CLUSTERED_COLUMN chart (2x2 table)")
        return ChartType.CLUSTERED_COLUMN.value

    # Single row → Horizontal bar (better for label readability)
    if n_rows == 1:
        logger.debug(f"  → HORIZONTAL_BAR chart (single row)")
        return ChartType.HORIZONTAL_BAR.value

    # Single column → Clustered column (standard comparison)
    if n_cols == 1:
        logger.debug(f"  → CLUSTERED_COLUMN chart (single column)")
        return ChartType.CLUSTERED_COLUMN.value

    # ==========================================================================
    # Priority 3: Dimension-based selection
    # ==========================================================================

    # More than 5 rows → Horizontal bar (better for many categories)
    if n_rows > 5:
        logger.debug(f"  → HORIZONTAL_BAR chart ({n_rows} rows > 5)")
        return ChartType.HORIZONTAL_BAR.value

    # More than 4 columns → 100% stacked column (part-to-whole)
    if n_cols > 4:
        logger.debug(f"  → STACKED_COLUMN_100 chart ({n_cols} cols > 4)")
        return ChartType.STACKED_COLUMN_100.value

    # ==========================================================================
    # Priority 4: Default → Clustered column
    # ==========================================================================
    logger.debug(f"  → CLUSTERED_COLUMN chart (default)")
    return ChartType.CLUSTERED_COLUMN.value


def get_xl_chart_type(chart_type: str) -> str:
    """
    Map internal chart type string to python-pptx XL_CHART_TYPE constant name.

    This function returns the string name of the XL_CHART_TYPE enum value
    that can be used with python-pptx when creating charts.

    Args:
        chart_type: Chart type string from select_chart_type() output

    Returns:
        String representation of XL_CHART_TYPE constant name

    Raises:
        ValueError: If chart_type is not recognized

    Examples:
        >>> get_xl_chart_type("clustered_column")
        'COLUMN_CLUSTERED'
        >>> get_xl_chart_type("horizontal_bar")
        'BAR_CLUSTERED'
    """
    mapping = {
        ChartType.CLUSTERED_COLUMN.value: "COLUMN_CLUSTERED",
        ChartType.HORIZONTAL_BAR.value: "BAR_CLUSTERED",
        ChartType.STACKED_COLUMN_100.value: "COLUMN_STACKED_100",
        ChartType.LINE.value: "LINE",
    }

    if chart_type not in mapping:
        raise ValueError(
            f"Unknown chart type: {chart_type}. "
            f"Valid types: {list(mapping.keys())}"
        )

    return mapping[chart_type]


def get_chart_type_description(chart_type: str) -> Dict[str, Any]:
    """
    Get documentation for a chart type.

    Args:
        chart_type: Chart type string from select_chart_type() output

    Returns:
        Documentation dictionary with keys:
        - xl_chart_type: python-pptx XL_CHART_TYPE constant name
        - description: Human-readable description
        - best_for: List of use cases
        - readability: Readability notes
        - examples: List of example tables

    Raises:
        ValueError: If chart_type is not recognized

    Examples:
        >>> docs = get_chart_type_description("clustered_column")
        >>> print(docs["description"])
        Standard vertical bar chart with grouped bars
        >>> print(docs["xl_chart_type"])
        COLUMN_CLUSTERED
    """
    if chart_type not in CHART_TYPE_DOCUMENTATION:
        raise ValueError(
            f"Unknown chart type: {chart_type}. "
            f"Valid types: {list(CHART_TYPE_DOCUMENTATION.keys())}"
        )

    return CHART_TYPE_DOCUMENTATION[chart_type].copy()


# =============================================================================
# Step 21: Generate PowerPoint
# =============================================================================

@trace_node("Step 21: Generate PowerPoint")
def generate_powerpoint_node(state: WorkflowState) -> dict:
    """
    Step 21: Generate PowerPoint presentation with charts from significant tables.

    This node creates a professional PowerPoint presentation where each significant
    table from Step 20 is rendered as an appropriate chart based on:
    - Table dimensions (rows x columns)
    - Semantic hints (if available)
    - Chart readability best practices

    Node Functionality:
    1. Load filtered_tables from state (Step 20 output)
    2. Get statistical_summary for test results
    3. Get metadata for variable labels
    4. Create PowerPoint presentation:
       - Add title slide with survey info
       - Add slide for each significant table with chart and data table
       - Add summary slide with key findings
    5. Save PowerPoint file
    6. Update state with powerpoint_file path

    Args:
        state: Current workflow state. Must contain:
            - filtered_tables: Significant tables from Step 20
            - statistical_summary: Statistical test results (optional)
            - new_metadata: Metadata for variable labels (optional)
            - raw_data: Original survey data for sample size (optional)

    Returns:
        Updated workflow state with:
            - powerpoint_file: Path to generated .pptx file
            - current_step: Set to STEP_21_GENERATE_POWERPOINT
            - errors: List of errors (appended if any occur)
            - warnings: List of warnings (appended if any occur)

    Example:
        >>> state = {
        ...     "filtered_tables": {
        ...         "tables": [
        ...             {
        ...                 "table_name": "gender_x_satisfaction",
        ...                 "row_variable": "gender",
        ...                 "column_variable": "satisfaction",
        ...                 "data": {
        ...                     "row_labels": ["Male", "Female"],
        ...                     "column_labels": ["Satisfied", "Neutral"],
        ...                     "counts": [[45, 32], [52, 28]]
        ...                 }
        ...             }
        ...         ]
        ...     },
        ...     "statistical_summary": {
        ...         "tables": [
        ...             {
        ...                 "table_name": "gender_x_satisfaction",
        ...                 "chi_square": 4.52,
        ...                 "p_value": 0.0335,
        ...                 "cramers_v": 0.18,
        ...                 "interpretation": "small"
        ...             }
        ...         ]
        ...     }
        ... }
        >>> new_state = generate_powerpoint_node(state)
        >>> print(new_state["powerpoint_file"])
        output/survey_analysis.pptx
    """
    logger.info("Step 21: Generating PowerPoint presentation")

    # Get required inputs from state
    filtered_tables = state.get("filtered_tables")
    statistical_summary = state.get("statistical_summary")
    new_metadata = state.get("new_metadata")
    raw_data = state.get("raw_data")
    original_metadata = state.get("original_metadata")
    config = state.get("config", DEFAULT_CONFIG)

    # Validate required inputs
    if not filtered_tables:
        error_msg = "No filtered_tables available in state. Cannot generate PowerPoint."
        logger.error(error_msg)
        return {
            "current_step": STEP_21_GENERATE_POWERPOINT,
            "errors": [error_msg],
        }

    try:
        # Import python-pptx
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.chart.data import CategoryChartData

        # Get tables from filtered_tables
        tables = filtered_tables.get("tables", [])
        logger.info(f"Found {len(tables)} significant tables for PowerPoint generation")

        # Build statistics lookup table
        stats_lookup = {}
        if statistical_summary:
            for table_stats in statistical_summary.get("tables", []):
                table_name = table_stats.get("table_name", "")
                stats_lookup[table_name] = table_stats

        # Get sample size and input file name
        sample_size = 0
        if raw_data is not None:
            sample_size = len(raw_data)
        input_file_name = "Unknown Survey"
        if original_metadata:
            input_file_name = original_metadata.get("file_name", "Unknown Survey")

        # Create output directory
        output_dir = Path(config.get("output_dir", "output"))
        output_dir.mkdir(parents=True, exist_ok=True)

        # Create PowerPoint presentation
        prs = Presentation()

        # ======================================================================
        # Add Title Slide
        # ======================================================================
        _add_title_slide(
            prs,
            title="Survey Analysis Results",
            subtitle=input_file_name,
            sample_size=sample_size
        )

        # ======================================================================
        # Add Slide for Each Significant Table
        # ======================================================================
        for table in tables:
            table_name = table.get("table_name", "unknown")
            row_variable = table.get("row_variable", "Unknown")
            column_variable = table.get("column_variable", "Unknown")

            logger.info(f"Creating slide for: {table_name}")

            # Get statistics for this table
            table_stats = stats_lookup.get(table_name, {})

            # Get variable labels from metadata
            row_label = _get_variable_label(row_variable, new_metadata)
            col_label = _get_variable_label(column_variable, new_metadata)

            # Create table slide
            try:
                _add_table_slide(
                    prs,
                    table=table,
                    row_label=row_label,
                    col_label=col_label,
                    table_stats=table_stats,
                    new_metadata=new_metadata
                )
            except Exception as e:
                logger.warning(f"Could not create slide for {table_name}: {e}")
                # Continue with next table

        # ======================================================================
        # Add Summary Slide
        # ======================================================================
        _add_summary_slide(
            prs,
            tables=tables,
            stats_lookup=stats_lookup,
            filtered_tables_summary=filtered_tables.get("summary", {})
        )

        # ======================================================================
        # Save PowerPoint File
        # ======================================================================
        powerpoint_path = output_dir / "survey_analysis.pptx"
        prs.save(str(powerpoint_path))
        logger.info(f"PowerPoint saved to: {powerpoint_path}")

        # ======================================================================
        # Return new state
        # ======================================================================
        return {
            "current_step": STEP_21_GENERATE_POWERPOINT,
            "powerpoint_file": str(powerpoint_path),
        }

    except ImportError as e:
        error_msg = f"Required library not found: {e}. Please install: pip install python-pptx"
        logger.error(error_msg)
        return {
            "current_step": STEP_21_GENERATE_POWERPOINT,
            "errors": [error_msg],
        }
    except Exception as e:
        error_msg = f"Unexpected error generating PowerPoint: {str(e)}"
        logger.error(error_msg, exc_info=True)
        return {
            "current_step": STEP_21_GENERATE_POWERPOINT,
            "errors": [error_msg],
        }


# =============================================================================
# Slide Creation Helper Functions
# =============================================================================

def _add_title_slide(
    prs,
    title: str,
    subtitle: str,
    sample_size: int
) -> None:
    """
    Add a title slide to the presentation.

    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Subtitle (file name)
        sample_size: Number of respondents
    """
    from pptx.util import Inches, Pt

    # Use blank layout for custom title slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Set slide background (optional - using default white)

    # Add title
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title

    # Format title - use professional font settings
    title_para = title_frame.paragraphs[0]
    title_para.font.name = TITLE_FONT["name"]
    title_para.font.size = Pt(TITLE_FONT["size"])
    title_para.font.bold = TITLE_FONT["bold"]
    title_para.font.color.rgb = get_market_research_color("primary")
    title_para.alignment = 1  # Center

    # Add subtitle
    top = Inches(2.7)
    subtitle_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle

    # Format subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.name = HEADING_FONT["name"]
    subtitle_para.font.size = Pt(HEADING_FONT["size"])
    subtitle_para.font.bold = HEADING_FONT["bold"]
    subtitle_para.font.color.rgb = get_market_research_color("secondary")
    subtitle_para.alignment = 1  # Center

    # Add date and sample size info
    top = Inches(4)
    info_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    info_frame = info_box.text_frame

    current_date = datetime.now().strftime("%B %d, %Y")
    info_frame.text = f"{current_date}  |  {sample_size:,} respondents"

    # Format info
    info_para = info_frame.paragraphs[0]
    info_para.font.name = BODY_FONT["name"]
    info_para.font.size = Pt(BODY_FONT["size"])
    info_para.font.color.rgb = get_market_research_color("muted")
    info_para.alignment = 1  # Center

    logger.debug("Title slide added")


def _add_table_slide(
    prs,
    table: Dict[str, Any],
    row_label: str,
    col_label: str,
    table_stats: Dict[str, Any],
    new_metadata: Optional[Dict[str, Any]]
) -> None:
    """
    Add a slide with chart and data table for a single cross-tabulation.

    Args:
        prs: Presentation object
        table: Table data dictionary
        row_label: Row variable display label
        col_label: Column variable display label
        table_stats: Statistical test results
        new_metadata: Metadata for value labels
    """
    from pptx.util import Inches, Pt
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData

    # Use blank layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Get table data
    table_data = table.get("data", {})
    row_labels = table_data.get("row_labels", [])
    column_labels = table_data.get("column_labels", [])
    counts = table_data.get("counts", [])

    # Get statistics
    chi_square = table_stats.get("chi_square", 0)
    p_value = table_stats.get("p_value", 1.0)
    cramers_v = table_stats.get("cramers_v", 0.0)
    interpretation = table_stats.get("interpretation", "unknown")
    is_significant = table_stats.get("is_significant", False)

    # =========================================================================
    # Add Title
    # =========================================================================
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = f"{row_label} × {col_label}"

    # Format title with professional font
    title_para = title_frame.paragraphs[0]
    title_para.font.name = TITLE_FONT["name"]
    title_para.font.size = Pt(TITLE_FONT["size"])
    title_para.font.bold = TITLE_FONT["bold"]
    title_para.font.color.rgb = get_market_research_color("primary")

    # =========================================================================
    # Select Chart Type
    # =========================================================================
    chart_type = select_chart_type(table)
    xl_chart_type_name = get_xl_chart_type(chart_type)
    xl_chart_type = getattr(XL_CHART_TYPE, xl_chart_type_name)

    logger.debug(f"  Chart type: {chart_type} ({xl_chart_type_name})")

    # =========================================================================
    # Add Chart
    # =========================================================================
    chart = None
    try:
        # Prepare chart data
        chart_data = CategoryChartData()

        # Set categories (row labels)
        # Use value labels if available
        display_row_labels = _get_value_labels(
            row_labels, table.get("row_variable"), new_metadata
        )
        chart_data.categories = display_row_labels[:10]  # Limit to 10 rows

        # Add series for each column
        display_col_labels = _get_value_labels(
            column_labels, table.get("column_variable"), new_metadata
        )

        # Transpose counts for python-pptx (series per column)
        if counts:
            n_rows = len(counts)
            n_cols = len(counts[0]) if n_rows > 0 else 0

            for col_idx in range(min(n_cols, 6)):  # Limit to 6 series
                col_name = display_col_labels[col_idx] if col_idx < len(display_col_labels) else f"Series {col_idx + 1}"
                series_values = [counts[row_idx][col_idx] if col_idx < len(counts[row_idx]) else 0 for row_idx in range(min(n_rows, 10))]
                chart_data.add_series(col_name, series_values)

        # Position chart (left side of slide)
        chart_left = Inches(0.5)
        chart_top = Inches(1.2)
        chart_width = Inches(4.5)
        chart_height = Inches(4)

        # Add chart to slide
        chart_shape = slide.shapes.add_chart(
            xl_chart_type, chart_left, chart_top, chart_width, chart_height, chart_data
        )
        chart = chart_shape.chart

        logger.debug(f"  Chart added: {xl_chart_type_name}")

    except Exception as e:
        logger.warning(f"  Could not add chart: {e}")

    # Apply styling to chart if it was created successfully
    if chart:
        try:
            _apply_chart_styling(chart, is_significant=is_significant)
        except Exception as e:
            logger.debug(f"  Could not apply chart styling: {e}")

    # =========================================================================
    # Add Data Table (right side of slide)
    # =========================================================================
    try:
        _add_data_table_to_slide(
            slide,
            row_labels=display_row_labels,
            col_labels=display_col_labels,
            counts=counts,
            x=Inches(5.2),
            y=Inches(1.2),
            width=Inches(4.5)
        )
    except Exception as e:
        logger.warning(f"  Could not add data table: {e}")

    # =========================================================================
    # Add Statistics Footer
    # =========================================================================
    left = Inches(0.5)
    top = Inches(5.4)
    width = Inches(9)
    height = Inches(0.4)

    footer_box = slide.shapes.add_textbox(left, top, width, height)
    footer_frame = footer_box.text_frame

    footer_text = f"χ² = {chi_square:.2f}, p = {p_value:.4f}, V = {cramers_v:.2f} ({interpretation})"
    footer_frame.text = footer_text

    # Format footer with professional font settings
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.name = FOOTER_FONT["name"]
    footer_para.font.size = Pt(FOOTER_FONT["size"])
    footer_para.font.italic = FOOTER_FONT["italic"]
    footer_para.font.color.rgb = get_market_research_color("muted")

    logger.debug(f"Table slide added: {row_label} × {col_label}")


def _add_summary_slide(
    prs,
    tables: List[Dict[str, Any]],
    stats_lookup: Dict[str, Dict[str, Any]],
    filtered_tables_summary: Dict[str, Any]
) -> None:
    """
    Add a summary slide with key findings.

    Args:
        prs: Presentation object
        tables: List of significant tables
        stats_lookup: Dictionary mapping table names to statistics
        filtered_tables_summary: Summary from filtering step
    """
    from pptx.util import Inches, Pt

    # Use blank layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # =========================================================================
    # Add Title
    # =========================================================================
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Summary of Key Findings"

    # Format title with professional font
    title_para = title_frame.paragraphs[0]
    title_para.font.name = TITLE_FONT["name"]
    title_para.font.size = Pt(TITLE_FONT["size"])
    title_para.font.bold = TITLE_FONT["bold"]
    title_para.font.color.rgb = get_market_research_color("primary")

    # =========================================================================
    # Add Summary Statistics
    # =========================================================================
    top = Inches(1.2)
    summary_box = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    summary_frame = summary_box.text_frame
    summary_frame.word_wrap = True

    total_analyzed = filtered_tables_summary.get("original_count", len(tables))
    significant_count = filtered_tables_summary.get("filtered_count", len(tables))

    # Build summary text
    summary_text = (
        f"Total tables analyzed: {total_analyzed}\n"
        f"Significant findings: {significant_count}"
    )

    if significant_count > 0:
        inclusion_rate = filtered_tables_summary.get("inclusion_rate", 0)
        summary_text += f"\nInclusion rate: {inclusion_rate:.1f}%"

    summary_frame.text = summary_text

    # Format summary with professional fonts
    for para in summary_frame.paragraphs:
        para.font.name = BODY_FONT["name"]
        para.font.size = Pt(BODY_FONT["size"])
        para.font.color.rgb = get_market_research_color("secondary")
        para.space_after = Pt(6)

    # =========================================================================
    # Add Top 3 Strongest Associations (by Cramer's V)
    # =========================================================================
    if tables:
        # Sort tables by Cramer's V
        sorted_tables = sorted(
            tables,
            key=lambda t: stats_lookup.get(t.get("table_name", ""), {}).get("cramers_v", 0),
            reverse=True
        )
        top_3 = sorted_tables[:3]

        if top_3:
            top = Inches(3)
            top_box = slide.shapes.add_textbox(left, top, width, Inches(2))
            top_frame = top_box.text_frame
            top_frame.word_wrap = True

            # Add heading
            p = top_frame.paragraphs[0]
            p.text = "Top 3 Strongest Associations:"
            p.font.name = HEADING_FONT["name"]
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = get_market_research_color("primary")
            p.space_after = Pt(12)

            # Add each table
            for i, table in enumerate(top_3):
                table_name = table.get("table_name", "Unknown")
                stats = stats_lookup.get(table_name, {})
                v = stats.get("cramers_v", 0)
                interp = stats.get("interpretation", "unknown")

                p = top_frame.add_paragraph()
                p.text = f"{i+1}. {table_name} (V = {v:.2f}, {interp})"
                p.font.name = BODY_FONT["name"]
                p.font.size = Pt(18)
                p.font.color.rgb = get_market_research_color("text")
                p.space_after = Pt(6)
                p.level = 0

    logger.debug("Summary slide added")


def _add_data_table_to_slide(
    slide,
    row_labels: List[str],
    col_labels: List[str],
    counts: List[List[int]],
    x,
    y,
    width
) -> None:
    """
    Add a formatted data table to a slide.

    Args:
        slide: Slide object
        row_labels: Row category labels
        col_labels: Column category labels
        counts: 2D list of counts
        x: Left position (Inches)
        y: Top position (Inches)
        width: Table width (Inches)
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_ANCHOR

    # Limit table size for readability
    max_rows = 8
    max_cols = 5

    num_rows = min(len(row_labels), max_rows) + 1  # +1 for header row
    num_cols = min(len(col_labels), max_cols) + 1  # +1 for row labels column

    # Calculate cell dimensions
    row_height = Inches(0.35)
    col_width = width / num_cols

    # Add table shape
    table_shape = slide.shapes.add_table(
        num_rows,
        num_cols,
        x, y, width, num_rows * row_height
    ).table

    # Format table
    # =========================================================================
    # Header row (column labels)
    # =========================================================================
    # First cell is empty or "Total"
    cell = table_shape.cell(0, 0)
    cell.text = ""
    _format_table_cell(cell, is_header=True)

    # Column labels
    for col_idx in range(min(len(col_labels), max_cols)):
        cell = table_shape.cell(0, col_idx + 1)
        cell.text = str(col_labels[col_idx])[:20]  # Truncate long labels
        _format_table_cell(cell, is_header=True)

    # Data rows
    for row_idx in range(min(len(row_labels), max_rows)):
        # Row label
        cell = table_shape.cell(row_idx + 1, 0)
        cell.text = str(row_labels[row_idx])[:25]  # Truncate long labels
        _format_table_cell(cell, is_header=False)

        # Counts
        if counts and row_idx < len(counts):
            for col_idx in range(min(len(col_labels), max_cols)):
                if col_idx < len(counts[row_idx]):
                    cell = table_shape.cell(row_idx + 1, col_idx + 1)
                    cell.text = str(counts[row_idx][col_idx])
                    _format_table_cell(cell, is_header=False)
                else:
                    cell = table_shape.cell(row_idx + 1, col_idx + 1)
                    cell.text = "0"
                    _format_table_cell(cell, is_header=False)


def _format_table_cell(cell, is_header: bool) -> None:
    """
    Format a table cell with appropriate styling.

    Args:
        cell: Table cell object
        is_header: Whether this is a header cell
    """
    from pptx.util import Pt

    # Get font settings from styling module
    if is_header:
        font_settings = TABLE_HEADER_FONT
        color_name = "primary"
    else:
        font_settings = TABLE_CELL_FONT
        color_name = "text"

    # Set font with professional settings
    para = cell.text_frame.paragraphs[0]
    para.font.name = font_settings["name"]
    para.font.size = Pt(font_settings["size"])

    if font_settings.get("bold"):
        para.font.bold = True
    if font_settings.get("italic"):
        para.font.italic = True

    para.font.color.rgb = get_market_research_color(color_name)

    # Set alignment
    para.alignment = 1  # Center


def _get_variable_label(variable_name: str, metadata: Optional[Dict[str, Any]]) -> str:
    """
    Get display label for a variable from metadata.

    Args:
        variable_name: Internal variable name
        metadata: Variable-centered metadata

    Returns:
        Display label (variable name if not found)
    """
    if not metadata:
        return variable_name

    variable_data = metadata.get(variable_name, {})
    label = variable_data.get("label", variable_name)

    return label if label else variable_name


def _get_value_labels(
    value_labels: List[str],
    variable_name: Optional[str],
    metadata: Optional[Dict[str, Any]]
) -> List[str]:
    """
    Get display labels for values from metadata.

    Args:
        value_labels: Internal value codes
        variable_name: Variable name for lookup
        metadata: Variable-centered metadata

    Returns:
        List of display labels
    """
    if not metadata or not variable_name:
        return value_labels

    variable_data = metadata.get(variable_name, {})
    value_labels_dict = variable_data.get("value_labels", {})

    # Map codes to labels
    display_labels = []
    for code in value_labels:
        label = value_labels_dict.get(str(code), str(code))
        display_labels.append(label)

    return display_labels


def _apply_chart_styling(chart, is_significant: bool = True) -> None:
    """
    Apply professional styling to a PowerPoint chart.

    This function adds color, formatting, and accessibility improvements
    to charts while keeping them natively editable in PowerPoint.

    Args:
        chart: PowerPoint chart object
        is_significant: Whether the chart shows significant results
    """
    from pptx.util import Pt
    from pptx.util import RGBColor

    # Style chart title
    if hasattr(chart, 'chart_title') and chart.chart_title:
        try:
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(20)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = (
                get_market_research_color("primary")
            )
        except Exception as e:
            logger.debug(f"Could not style chart title: {e}")

    # Style series colors
    try:
        for idx, series in enumerate(chart.series):
            # Get color from palette
            color_hex = get_chart_color(idx, hex_format=True)
            color_rgb = hex_to_rgb(color_hex)

            # Apply fill color
            series.format.fill.fore_color.rgb = RGBColor(*color_rgb)

            # Add subtle border for definition
            series.format.line.color.rgb = RGBColor(*color_rgb)
            series.format.line.width = Pt(1)
    except Exception as e:
        logger.debug(f"Could not style chart series: {e}")

    # Style chart area (background)
    try:
        chart.chart_area.format.fill.fore_color.rgb = RGBColor(255, 255, 255)
    except Exception as e:
        logger.debug(f"Could not style chart area: {e}")

    logger.debug("Chart styling applied")


# =============================================================================
# Utility Functions
# =============================================================================

def validate_table_dimensions(
    n_rows: int,
    n_cols: int,
    max_rows: int = 20,
    max_cols: int = 10
) -> Tuple[bool, Optional[str]]:
    """
    Validate table dimensions for chart generation.

    Args:
        n_rows: Number of rows in table
        n_cols: Number of columns in table
        max_rows: Maximum recommended rows (default: 20)
        max_cols: Maximum recommended columns (default: 10)

    Returns:
        Tuple of (is_valid, warning_message):
        - is_valid: True if dimensions are within bounds
        - warning_message: Warning if dimensions exceed recommendations, None otherwise

    Examples:
        >>> validate_table_dimensions(5, 3)
        (True, None)
        >>> validate_table_dimensions(25, 3)
        (False, 'Table has 25 rows, exceeding recommended maximum of 20')
        >>> validate_table_dimensions(5, 12)
        (False, 'Table has 12 columns, exceeding recommended maximum of 10')
    """
    if n_rows > max_rows:
        return False, f"Table has {n_rows} rows, exceeding recommended maximum of {max_rows}"

    if n_cols > max_cols:
        return False, f"Table has {n_cols} columns, exceeding recommended maximum of {max_cols}"

    return True, None

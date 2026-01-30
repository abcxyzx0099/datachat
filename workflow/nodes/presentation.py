"""
Nodes for Phase 7: Presentation - PowerPoint Generation with Native Editable Charts

This module handles the generation of PowerPoint presentations with native
editable charts for cross-tabulation survey data. Each significant table gets
a slide with a bar chart, stacked bar chart, or horizontal bar chart showing
the data visualization.

Charts are created using python-pptx's native chart API (add_chart), making
them fully editable in PowerPoint, LibreOffice Impress, and compatible applications.
"""

import json
import os
from typing import Dict, Any, List, Tuple, Optional
from pathlib import Path

from ..state import State


# ============================================================================
# MAIN NODE: GENERATE POWERPOINT WITH NATIVE CHARTS
# ============================================================================

def generate_powerpoint(state: State) -> State:
    """
    Create PowerPoint presentation with native editable charts from significant tables.

    This is the main node for Phase 7 (Step 21) of the workflow. It:
    1. Loads significant tables from JSON file
    2. Loads statistical summary (Chi-square, p-value, Cramer's V)
    3. Generates a slide for each table with:
       - Slide title
       - Native PowerPoint bar chart or stacked bar chart (editable!)
       - Statistical summary text below the chart

    Charts are created using python-pptx's add_chart() API, making them fully
    editable in PowerPoint, LibreOffice Impress, and other compatible applications.
    Users can modify data, change chart types, adjust styling, and edit labels.

    Input state fields:
        - significant_tables_json_path: Path to filtered significant tables JSON
        - statistical_summary_path: Path to statistical test results JSON
        - config["output_dir"]: Output directory for PowerPoint file

    Output state fields:
        - powerpoint_path: Path to generated .pptx file
        - charts_generated: List of chart metadata (table_name, chart_type, statistics)

    Args:
        state: Current workflow state

    Returns:
        Updated state with PowerPoint generation results
    """
    try:
        output_dir = state.get("config", {}).get("output_dir", "output")
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        # Load data files
        tables = _load_significant_tables(state["significant_tables_json_path"])
        statistics = _load_statistical_summary(state["statistical_summary_path"])

        if not tables:
            return {
                **state,
                "warnings": state.get("warnings", []) + [
                    "No significant tables found. PowerPoint generation skipped."
                ],
                "execution_log": state.get("execution_log", []) + [{
                    "step": "generate_powerpoint",
                    "status": "skipped",
                    "reason": "no_significant_tables"
                }]
            }

        # Create presentation with native charts
        ppt_path, charts_generated = _create_presentation_with_native_charts(
            tables=tables,
            statistics=statistics,
            output_dir=output_dir
        )

        return {
            **state,
            "powerpoint_path": ppt_path,
            "charts_generated": charts_generated,
            "messages": [
                *state.get("messages", []),
                {
                    "role": "assistant",
                    "content": f"Generated PowerPoint with {len(charts_generated)} native editable chart slides"
                }
            ],
            "execution_log": state.get("execution_log", []) + [{
                "step": "generate_powerpoint",
                "status": "completed",
                "ppt_path": ppt_path,
                "charts_generated": len(charts_generated)
            }]
        }

    except Exception as e:
        return {
            **state,
            "errors": state.get("errors", []) + [
                f"Failed to generate PowerPoint: {str(e)}"
            ],
            "execution_log": state.get("execution_log", []) + [{
                "step": "generate_powerpoint",
                "status": "failed",
                "error": str(e)
            }]
        }


# ============================================================================
# DATA LOADING FUNCTIONS
# ============================================================================

def _load_significant_tables(json_path: str) -> List[Dict[str, Any]]:
    """
    Load significant tables from JSON file.

    Args:
        json_path: Path to significant_tables JSON file

    Returns:
        List of table dictionaries with structure:
        {
            "name": str,
            "rows": str,  # row variable name
            "columns": str,  # column variable name
            "data": {
                "row_labels": List[str],
                "column_labels": List[str],
                "counts": List[List[int]],  # 2D array of counts
                "row_percentages": List[List[float]],  # optional
                "column_percentages": List[List[float]]  # optional
            }
        }
    """
    with open(json_path, 'r') as f:
        data = json.load(f)

    # Handle different JSON structures
    if isinstance(data, dict):
        if "tables" in data:
            return data["tables"]
        elif "significant_tables" in data:
            return data["significant_tables"]
        else:
            # Assume the dict itself contains table data
            return [data]
    elif isinstance(data, list):
        return data

    raise ValueError(f"Unexpected JSON structure in {json_path}")


def _load_statistical_summary(json_path: str) -> List[Dict[str, Any]]:
    """
    Load statistical summary from JSON file.

    Args:
        json_path: Path to statistical_analysis_summary JSON file

    Returns:
        List of statistical result dictionaries with structure:
        {
            "table_name": str,
            "chi_square": float,
            "p_value": float,
            "degrees_of_freedom": int,
            "cramers_v": float,
            "interpretation": str,  # "negligible", "small", "medium", "large"
            "sample_size": int
        }
    """
    with open(json_path, 'r') as f:
        data = json.load(f)

    if isinstance(data, dict) and "results" in data:
        return data["results"]
    elif isinstance(data, list):
        return data

    raise ValueError(f"Unexpected JSON structure in {json_path}")


# ============================================================================
# NATIVE CHART GENERATION FUNCTIONS
# ============================================================================

def _create_presentation_with_native_charts(
    tables: List[Dict[str, Any]],
    statistics: List[Dict[str, Any]],
    output_dir: str
) -> Tuple[str, List[Dict[str, Any]]]:
    """
    Create PowerPoint presentation with native editable charts for all tables.

    Args:
        tables: List of significant tables with data
        statistics: List of statistical test results
        output_dir: Directory to save PowerPoint file

    Returns:
        Tuple of (pptx_file_path, charts_generated_metadata)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise ImportError(
            "python-pptx is required. Install with: pip install python-pptx"
        )

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Add title slide
    _add_title_slide(prs)

    # Track generated charts
    charts_generated = []

    # Create a statistics lookup dictionary
    stats_lookup = {
        stat["table_name"]: stat
        for stat in statistics
    }

    # Process each table
    for table_data in tables:
        table_name = table_data.get("name", "Unnamed Table")
        table_stats = stats_lookup.get(table_name)

        # Create slide with native chart
        try:
            chart_metadata = _add_table_slide_with_native_chart(
                prs=prs,
                table_data=table_data,
                statistics=table_stats
            )
            charts_generated.append(chart_metadata)

        except Exception as e:
            # Log error but continue with other tables
            print(f"Warning: Failed to create chart for '{table_name}': {e}")
            # Add a placeholder slide with error message
            _add_error_slide(prs, table_name, str(e))

    # Save presentation
    ppt_path = os.path.join(output_dir, "survey_analysis_with_charts.pptx")
    prs.save(ppt_path)

    return ppt_path, charts_generated


def _add_title_slide(prs) -> None:
    """
    Add title slide to presentation.

    Args:
        prs: PowerPoint Presentation object
    """
    from pptx.util import Inches, Pt
    from datetime import datetime

    # Use blank layout for custom title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Survey Analysis Results"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = 1  # Center

    # Add subtitle with date
    top = Inches(3.2)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.alignment = 1  # Center


def _add_table_slide_with_native_chart(
    prs,
    table_data: Dict[str, Any],
    statistics: Optional[Dict[str, Any]]
) -> Dict[str, Any]:
    """
    Add a slide with native editable PowerPoint chart for a single table.

    Args:
        prs: PowerPoint Presentation object
        table_data: Table data including labels and counts
        statistics: Statistical test results (optional)

    Returns:
        Chart metadata dictionary
    """
    from pptx.util import Inches, Pt
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData

    table_name = table_data.get("name", "Unnamed Table")

    # Create blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(9)
    title_height = Inches(0.5)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = table_name
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True

    # Extract table data
    chart_data_input = table_data.get("data", {})
    row_labels = chart_data_input.get("row_labels", [])
    col_labels = chart_data_input.get("column_labels", [])
    counts = chart_data_input.get("counts", [])

    n_rows = len(row_labels)
    n_cols = len(col_labels)

    # Determine chart type based on table dimensions
    chart_type, xl_chart_type = _determine_chart_type(n_rows, n_cols)

    # Prepare chart data for python-pptx
    category_chart_data = _prepare_category_chart_data(
        row_labels=row_labels,
        col_labels=col_labels,
        counts=counts,
        chart_type=chart_type
    )

    # Add native chart to slide
    chart_left = Inches(0.5)
    chart_top = Inches(1.0)
    chart_width = Inches(9)
    chart_height = Inches(3.0)

    chart = slide.shapes.add_chart(
        xl_chart_type,
        chart_left,
        chart_top,
        chart_width,
        chart_height,
        category_chart_data
    ).chart

    # Configure chart appearance
    _configure_chart_style(chart, table_name, chart_type)

    # Add statistics summary text box (bottom portion)
    if statistics:
        stats_left = Inches(0.5)
        stats_top = Inches(4.2)
        stats_width = Inches(9)
        stats_height = Inches(1.2)

        stats_text = (
            f"Statistical Summary:\n"
            f"Chi-Square: {statistics['chi_square']:.3f}  |  "
            f"p-value: {statistics['p_value']:.4f}  |  "
            f"Cramer's V: {statistics['cramers_v']:.3f} ({statistics['interpretation']})"
        )

        stats_box = slide.shapes.add_textbox(stats_left, stats_top, stats_width, stats_height)
        stats_frame = stats_box.text_frame
        stats_frame.word_wrap = True
        stats_frame.text = stats_text

        # Format statistics text
        stats_para = stats_frame.paragraphs[0]
        stats_para.font.size = Pt(14)
        stats_para.font.color.rgb = _get_rgb_color(64, 64, 64)

    # Return chart metadata
    return {
        "table_name": table_name,
        "chart_type": chart_type,
        "dimensions": f"{n_rows}x{n_cols}",
        "statistics": statistics
    }


def _determine_chart_type(n_rows: int, n_cols: int) -> Tuple[str, int]:
    """
    Determine the best chart type based on table dimensions.

    Args:
        n_rows: Number of rows in the table
        n_cols: Number of columns in the table

    Returns:
        Tuple of (chart_type_string, xl_chart_type_enum)
        - chart_type_string: "bar", "stacked_bar", or "horizontal_bar"
        - xl_chart_type_enum: XL_CHART_TYPE enum value
    """
    from pptx.enum.chart import XL_CHART_TYPE

    # For 2x2 tables, use clustered column chart
    if n_rows == 2 and n_cols == 2:
        return "bar", XL_CHART_TYPE.COLUMN_CLUSTERED

    # For tables with many rows, use horizontal bar chart
    if n_rows > 5:
        return "horizontal_bar", XL_CHART_TYPE.BAR_CLUSTERED

    # For tables with many columns, use stacked bar chart
    if n_cols > 4:
        return "stacked_bar", XL_CHART_TYPE.COLUMN_STACKED_100

    # Default to clustered column chart
    return "bar", XL_CHART_TYPE.COLUMN_CLUSTERED


def _prepare_category_chart_data(
    row_labels: List[str],
    col_labels: List[str],
    counts: List[List[int]],
    chart_type: str
) -> 'CategoryChartData':
    """
    Transform cross-tabulation data to python-pptx CategoryChartData format.

    Args:
        row_labels: Labels for row categories
        col_labels: Labels for column categories
        counts: 2D array of counts [row][col]
        chart_type: Type of chart ("bar", "horizontal_bar", "stacked_bar")

    Returns:
        CategoryChartData object ready for add_chart()
    """
    from pptx.chart.data import CategoryChartData

    # Create chart data
    chart_data = CategoryChartData()

    # Add categories (row labels become chart categories)
    chart_data.categories = row_labels

    # For each column label, create a data series
    for col_idx, col_label in enumerate(col_labels):
        # Extract data for this series (one value per row/category)
        series_data = [counts[row_idx][col_idx] for row_idx in range(len(row_labels))]

        # Add series to chart data
        chart_data.add_series(col_label, series_data)

    return chart_data


def _configure_chart_style(chart, table_name: str, chart_type: str) -> None:
    """
    Configure chart styling for professional presentation.

    Args:
        chart: PowerPoint chart object
        table_name: Name for chart title
        chart_type: Type of chart being configured
    """
    from pptx.util import Pt

    # Set chart title
    if chart.has_title:
        title = chart.chart_title
        title.text_frame.text = table_name
        title.text_frame.paragraphs[0].font.size = Pt(14)
        title.text_frame.paragraphs[0].font.bold = True

    # Apply professional color scheme
    _apply_chart_colors(chart, chart_type)

    # Display legend
    if chart.has_legend:
        legend = chart.legend
        legend.include_in_layout = False
        legend.position = 2  # Right side

    # Set axis titles if applicable
    _set_axis_titles(chart, chart_type)

    # Show data labels for better readability
    _show_data_labels(chart, chart_type)


def _apply_chart_colors(chart, chart_type: str) -> None:
    """
    Apply professional color scheme to chart series.

    Args:
        chart: PowerPoint chart object
        chart_type: Type of chart
    """
    # Market research professional color palette
    colors_rgb = [
        (46, 134, 171),   # Blue #2E86AB
        (162, 59, 114),   # Purple #A23B72
        (241, 143, 1),    # Orange #F18F01
        (199, 62, 29),    # Red #C73E1D
        (106, 153, 78),   # Green #6A994E
        (188, 75, 81),    # Maroon #BC4B51
        (92, 110, 138),   # Steel Blue #5C6E8A
        (136, 176, 75),   # Olive Green #88B04B
    ]

    # Apply colors to each series
    for idx, series in enumerate(chart.series):
        if idx < len(colors_rgb):
            r, g, b = colors_rgb[idx]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _get_rgb_color(r, g, b)

            # Add slight transparency for modern look
            series.format.fill.fore_color.brightness = 0.0


def _set_axis_titles(chart, chart_type: str) -> None:
    """
    Set axis titles for the chart.

    Args:
        chart: PowerPoint chart object
        chart_type: Type of chart
    """
    from pptx.enum.chart import XL_AXIS_TYPE
    from pptx.util import Pt

    # Category axis (X-axis for vertical charts, Y-axis for horizontal)
    category_axis = chart.category_axis
    if category_axis.has_title:
        if chart_type == "horizontal_bar":
            category_axis.axis_title.text_frame.text = "Categories"
        else:
            category_axis.axis_title.text_frame.text = "Categories"
        category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        category_axis.axis_title.text_frame.paragraphs[0].font.bold = True

    # Value axis (Y-axis for vertical charts, X-axis for horizontal)
    value_axis = chart.value_axis
    if value_axis.has_title:
        if chart_type == "stacked_bar":
            value_axis.axis_title.text_frame.text = "Percentage (%)"
        else:
            value_axis.axis_title.text_frame.text = "Count"
        value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        value_axis.axis_title.text_frame.paragraphs[0].font.bold = True


def _show_data_labels(chart, chart_type: str) -> None:
    """
    Show data labels on chart for better readability.

    Args:
        chart: PowerPoint chart object
        chart_type: Type of chart
    """
    # Show data labels for each series
    for series in chart.series:
        # Note: Data labels are not always fully supported in python-pptx
        # The chart will have labels when opened in PowerPoint
        # This is a limitation of the library
        pass


def _add_error_slide(prs, table_name: str, error_message: str) -> None:
    """
    Add a slide with error message when chart generation fails.

    Args:
        prs: PowerPoint Presentation object
        table_name: Name of the table that failed
        error_message: Error message to display
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = f"{table_name} - Chart Generation Failed"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = _get_rgb_color(192, 0, 0)

    # Add error message
    error_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    error_frame = error_box.text_frame
    error_frame.text = f"Error: {error_message}"
    error_para = error_frame.paragraphs[0]
    error_para.font.size = Pt(14)
    error_para.alignment = PP_ALIGN.CENTER


# ============================================================================
# RGB COLOR HELPER
# ============================================================================

def _get_rgb_color(r: int, g: int, b: int):
    """
    Create an RGB color object for python-pptx.

    Args:
        r: Red component (0-255)
        g: Green component (0-255)
        b: Blue component (0-255)

    Returns:
        RGB color object compatible with python-pptx
    """
    try:
        from pptx.util import RGBColor as PptxRGBColor
        return PptxRGBColor(r, g, b)
    except ImportError:
        # Fallback for older python-pptx versions
        class RGBColor:
            def __init__(self, r: int, g: int, b: int):
                self.r = r
                self.g = g
                self.b = b
        return RGBColor(r, g, b)

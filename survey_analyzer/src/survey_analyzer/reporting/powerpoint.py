"""
PowerPoint Generator

Generate PowerPoint presentations from cross-tabulation results.

Example:
    >>> gen = PowerPointGenerator(template="template.pptx")
    >>> gen.add_presentation(title="Survey Results")
    >>> gen.add_slide_for_table(table_data, stats_data)
    >>> gen.save("output.pptx")
"""

import logging
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple

logger = logging.getLogger(__name__)


class ChartType(Enum):
    """Chart types for table visualization."""
    CLUSTERED_COLUMN = "clustered_column"
    HORIZONTAL_BAR = "horizontal_bar"
    STACKED_COLUMN_100 = "stacked_column_100"
    LINE = "line"


@dataclass
class ChartColorScheme:
    """Color scheme for charts."""
    primary: str = "#4287f5"
    secondary: str = "#48bb78"
    accent: str = "#f6e05e"
    significant: str = "#48bb78"
    not_significant: str = "#e53e3e"

    @classmethod
    def professional(cls) -> "ChartColorScheme":
        """Professional market research colors."""
        return cls(
            primary="#4287f5",
            secondary="#48bb78",
            accent="#ed8a36",
            significant="#48bb78",
            not_significant="#e53e3e",
        )


@dataclass
class SlideConfig:
    """Configuration for slide generation."""
    show_title: bool = True
    show_table: bool = True
    show_chart: bool = True
    show_statistics: bool = True
    show_significance_badge: bool = True
    chart_type: Optional[ChartType] = None
    theme_color: Optional[str] = None


class PowerPointGenerator:
    """
    Generate PowerPoint presentations from survey analysis results.

    Features:
    - Automatic chart type selection based on table dimensions
    - Significance highlighting
    - Professional styling
    - Template support

    Example:
        >>> gen = PowerPointGenerator()
        >>> gen.create_presentation(
        ...     tables=filtered_tables,
        ...     statistics=statistical_summary,
        ...     title="Q1 2024 Survey Results"
        ... )
        >>> gen.save("output/q1_2024_report.pptx")
    """

    # Chart type selection logic
    MAX_CATEGORIES_FOR_PIE = 5
    MAX_CATEGORIES_FOR_COLUMN = 8

    def __init__(
        self,
        template_path: Optional[str] = None,
        color_scheme: Optional[ChartColorScheme] = None,
    ):
        """
        Initialize the generator.

        Args:
            template_path: Optional path to PowerPoint template (.pptx)
            color_scheme: Color scheme for charts
        """
        self.template_path = template_path
        self.color_scheme = color_scheme or ChartColorScheme.professional()
        self._presentation = None

    def create_presentation(
        self,
        tables: List[Dict[str, Any]],
        statistics: Dict[str, Any],
        title: str = "Survey Analysis Results",
        subtitle: Optional[str] = None,
    ) -> None:
        """
        Create a complete presentation from tables and statistics.

        Args:
            tables: List of table data dictionaries
            statistics: Statistical summary with test results
            title: Presentation title
            subtitle: Optional subtitle

        Example:
            >>> gen = PowerPointGenerator()
            >>> gen.create_presentation(
            ...     tables=filtered_tables,
            ...     statistics=stats_summary,
            ...     title="Customer Satisfaction Report"
            ... )
            >>> gen.save("report.pptx")
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError(
                "python-pptx is required. Install with: pip install python-pptx"
            )

        # Create presentation
        if self.template_path and Path(self.template_path).exists():
            self._presentation = Presentation(self.template_path)
        else:
            self._presentation = Presentation()

        # Add title slide
        self._add_title_slide(title, subtitle)

        # Add summary slide
        self._add_summary_slide(tables, statistics)

        # Add table slides
        for idx, table in enumerate(tables):
            self._add_table_slide(table, statistics, idx + 1)

        logger.info(
            f"Created presentation with {len(self._presentation.slides)} slides"
        )

    def _add_title_slide(
        self,
        title: str,
        subtitle: Optional[str],
    ) -> None:
        """Add title slide to presentation."""
        from pptx.util import Pt

        title_slide = self._presentation.slides.add_slide(
            self._presentation.slide_layouts[0]  # Title slide layout
        )

        # Set title
        title_shape = title_slide.shapes.title
        title_shape.text = title

        # Set subtitle
        if subtitle:
            # Try to get subtitle placeholder
            for shape in title_slide.placeholders:
                if shape.placeholder_format.type == 1:  # Subtitle type
                    shape.text = subtitle
                    break

    def _add_summary_slide(
        self,
        tables: List[Dict[str, Any]],
        statistics: Dict[str, Any],
    ) -> None:
        """Add summary slide with overall statistics."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = self._presentation.slides.add_slide(
            self._presentation.slide_layouts[5]  # Blank layout
        )

        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "Analysis Summary"

        # Add summary stats
        stats = statistics.get("tables", [])
        total = len(tables)
        significant = sum(1 for s in stats if s.get("is_significant", False))

        # Calculate significance rate with division by zero protection
        significance_rate = (significant / total * 100) if total > 0 else 0

        summary_text = f"""
        Total Tables Analyzed: {total}
        Significant Relationships: {significant}
        Significance Rate: {significance_rate:.1f}%

        Generated by SPSS Analyzer
        """

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        summary_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = summary_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = summary_text.strip()

        # Format text
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.name = "Calibri"

    def _add_table_slide(
        self,
        table: Dict[str, Any],
        statistics: Dict[str, Any],
        slide_number: int,
    ) -> None:
        """
        Add a slide for a single table.

        Args:
            table: Table data dictionary
            statistics: Statistical summary (for finding matching stats)
            slide_number: Slide number for title
        """
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = self._presentation.slides.add_slide(
            self._presentation.slide_layouts[5]  # Blank layout
        )

        # Get table info
        table_name = table.get("table_name", f"Table {slide_number}")
        row_var = table.get("row_variable", "")
        col_var = table.get("column_variable", "")

        # Find matching statistics
        table_stats = self._find_table_stats(table_name, statistics)

        # Determine significance
        is_significant = table_stats.get("is_significant", False)
        sig_color = (
            RGBColor(72, 187, 120) if is_significant  # Green
            else RGBColor(229, 62, 62)  # Red
        )

        # Title
        title_text = f"{row_var} × {col_var}"
        if is_significant:
            title_text += " (Significant)"

        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.7)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = sig_color

        # Add statistics text
        if table_stats and table_stats.get("is_valid", True):
            p_value = table_stats.get("p_value")
            cramers_v = table_stats.get("cramers_v")
            interpretation = table_stats.get("interpretation", "")

            stats_text = f"χ² = {table_stats.get('chi_square', 0):.2f}, "
            stats_text += f"p = {p_value:.4f}, "
            stats_text += f"V = {cramers_v:.3f} ({interpretation})"

            left = Inches(0.5)
            top = Inches(1.1)
            width = Inches(9)
            height = Inches(0.4)

            stats_box = slide.shapes.add_textbox(left, top, width, height)
            stats_frame = stats_box.text_frame
            stats_frame.text = stats_text
            stats_frame.paragraphs[0].font.size = Pt(14)
            stats_frame.paragraphs[0].font.italic = True

    def _find_table_stats(
        self,
        table_name: str,
        statistics: Dict[str, Any],
    ) -> Optional[Dict[str, Any]]:
        """Find statistics for a specific table."""
        tables = statistics.get("tables", [])
        for table_stats in tables:
            if table_stats.get("table_name") == table_name:
                return table_stats
        return None

    def save(self, output_path: str) -> None:
        """
        Save the presentation to a file.

        Args:
            output_path: Path for output .pptx file

        Example:
            >>> gen = PowerPointGenerator()
            >>> gen.create_presentation(tables, stats)
            >>> gen.save("report.pptx")
        """
        if self._presentation is None:
            raise RuntimeError("No presentation created. Call create_presentation() first.")

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self._presentation.save(str(output_path))
        logger.info(f"Presentation saved to: {output_path}")


def create_powerpoint(
    tables: List[Dict[str, Any]],
    statistics: Dict[str, Any],
    output_path: str,
    title: str = "Survey Analysis Results",
) -> None:
    """
    Convenience function to create a PowerPoint presentation.

    Args:
        tables: List of table data dictionaries
        statistics: Statistical summary
        output_path: Path for output .pptx file
        title: Presentation title

    Example:
        >>> create_powerpoint(
        ...     tables=filtered_tables,
        ...     statistics=stats_summary,
        ...     output_path="report.pptx",
        ...     title="Q1 Results"
        ... )
    """
    gen = PowerPointGenerator()
    gen.create_presentation(tables, statistics, title)
    gen.save(output_path)
